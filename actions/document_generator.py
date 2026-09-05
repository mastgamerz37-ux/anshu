"""
actions/document_generator.py — Autonomous Document, Presentation & Spreadsheet Generation for ANSH

Supports:
- PowerPoint (.pptx) Presentations with custom slide layouts & styling
- Excel (.xlsx) Spreadsheets with formulas, colored headers & styled tables
- PDF (.pdf) Reports & Invoices
"""
from __future__ import annotations

import os
import sys
import json
import subprocess
from pathlib import Path
from datetime import datetime
from typing import Optional, List, Dict, Any

def _get_output_dir() -> Path:
    desktop = Path.home() / "Desktop"
    docs_dir = desktop / "ANSH Generated Documents"
    docs_dir.mkdir(parents=True, exist_ok=True)
    return docs_dir


def generate_presentation(
    topic: str,
    slides_data: Optional[List[Dict[str, Any]]] = None,
    num_slides: int = 5,
    player=None
) -> str:
    """
    Generate a PowerPoint (.pptx) presentation.
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        prs = Presentation()
        # Set 16:9 widescreen
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        output_dir = _get_output_dir()
        clean_topic = "".join(c for c in topic if c.isalnum() or c in (" ", "_", "-")).strip()
        filename = f"{clean_topic[:30].replace(' ', '_')}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        file_path = output_dir / filename

        # If slides data not provided, generate structured outline via Gemini
        if not slides_data:
            from core.task_llm import call_task_llm
            prompt = f"""Generate a structured {num_slides}-slide presentation on the topic: "{topic}".
Return ONLY a valid JSON array of objects with keys: "title" (string), "subtitle" (optional string), "bullets" (array of 3-5 concise strings).
No markdown, no backticks, only JSON."""
            raw = call_task_llm(prompt=prompt).strip()
            if raw.startswith("```"):
                raw = "\n".join(raw.split("\n")[1:-1]).strip()
            slides_data = json.loads(raw)

        # Title Slide
        blank_slide_layout = prs.slide_layouts[6]
        
        # Color Theme: Dark Cyber / Modern Blue
        bg_color = RGBColor(10, 15, 25)
        accent_color = RGBColor(0, 212, 255)
        text_color = RGBColor(240, 245, 255)
        subtext_color = RGBColor(160, 180, 205)

        # Slide 1: Title
        slide1 = prs.slides.add_slide(blank_slide_layout)
        txBox = slide1.shapes.add_textbox(Inches(1.5), Inches(2.2), Inches(10.3), Inches(3.0))
        tf = txBox.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = topic.upper()
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = accent_color
        p.alignment = PP_ALIGN.LEFT

        p2 = tf.add_paragraph()
        p2.text = f"Prepared by ANSH Autonomous AI System • {datetime.now().strftime('%B %Y')}"
        p2.font.size = Pt(18)
        p2.font.color.rgb = subtext_color
        p2.alignment = PP_ALIGN.LEFT

        # Content Slides
        for s_idx, slide_info in enumerate(slides_data, start=2):
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # Header
            header_box = slide.shapes.add_textbox(Inches(1.0), Inches(0.8), Inches(11.3), Inches(1.2))
            htf = header_box.text_frame
            hp = htf.paragraphs[0]
            hp.text = slide_info.get("title", f"Key Insights {s_idx - 1}")
            hp.font.size = Pt(32)
            hp.font.bold = True
            hp.font.color.rgb = accent_color

            # Bullets
            content_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.3), Inches(4.5))
            ctf = content_box.text_frame
            ctf.word_wrap = True

            bullets = slide_info.get("bullets", [])
            for b_idx, bullet_text in enumerate(bullets):
                bp = ctf.paragraphs[0] if b_idx == 0 else ctf.add_paragraph()
                bp.text = f"•  {bullet_text}"
                bp.font.size = Pt(20)
                bp.font.color.rgb = text_color
                bp.space_after = Pt(14)

        prs.save(str(file_path))

        # Open in default app
        if sys.platform == "win32":
            os.startfile(str(file_path))
        elif sys.platform == "darwin":
            subprocess.run(["open", str(file_path)])
        else:
            subprocess.run(["xdg-open", str(file_path)])

        if player:
            player.write_log(f"[DocGen] 📊 Generated PPTX: {filename}")

        return f"Successfully generated and launched PowerPoint presentation: '{filename}' ({len(slides_data) + 1} slides).\nSaved to: {file_path}"

    except Exception as e:
        print(f"[DocGen] PPTX error: {e}")
        return f"Failed to generate presentation: {e}"


def generate_spreadsheet(
    topic: str,
    table_data: Optional[Dict[str, Any]] = None,
    player=None
) -> str:
    """
    Generate an Excel (.xlsx) spreadsheet with formatted headers and data.
    """
    try:
        import openpyxl
        from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
        from openpyxl.utils import get_column_letter

        output_dir = _get_output_dir()
        clean_topic = "".join(c for c in topic if c.isalnum() or c in (" ", "_", "-")).strip()
        filename = f"{clean_topic[:30].replace(' ', '_')}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.xlsx"
        file_path = output_dir / filename

        # If data not provided, generate structured table via Gemini
        if not table_data:
            from core.task_llm import call_task_llm
            prompt = f"""Generate a structured spreadsheet table for: "{topic}".
Return ONLY a valid JSON object with keys:
- "title": (string)
- "headers": (array of strings, e.g. ["Category", "Item", "Cost", "Quantity", "Total"])
- "rows": (array of arrays containing data values matching the headers)
No markdown, no backticks, only JSON."""
            raw = call_task_llm(prompt=prompt).strip()
            if raw.startswith("```"):
                raw = "\n".join(raw.split("\n")[1:-1]).strip()
            table_data = json.loads(raw)

        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = (table_data.get("title") or "Sheet1")[:30]

        # Styling
        header_fill = PatternFill(start_color="0A2540", end_color="0A2540", fill_type="solid")
        header_font = Font(name="Segoe UI", size=11, bold=True, color="00D4FF")
        data_font = Font(name="Segoe UI", size=10)
        border_thin = Border(
            left=Side(style='thin', color='D0D7DE'),
            right=Side(style='thin', color='D0D7DE'),
            top=Side(style='thin', color='D0D7DE'),
            bottom=Side(style='thin', color='D0D7DE')
        )

        headers = table_data.get("headers", ["Column 1", "Column 2", "Column 3"])
        rows = table_data.get("rows", [])

        # Write Headers
        for col_num, header in enumerate(headers, 1):
            cell = ws.cell(row=1, column=col_num, value=header)
            cell.fill = header_fill
            cell.font = header_font
            cell.alignment = Alignment(horizontal="center", vertical="center")
            cell.border = border_thin

        # Write Rows
        for row_num, row_data in enumerate(rows, 2):
            for col_num, val in enumerate(row_data, 1):
                cell = ws.cell(row=row_num, column=col_num, value=val)
                cell.font = data_font
                cell.border = border_thin
                if isinstance(val, (int, float)):
                    cell.alignment = Alignment(horizontal="right")
                else:
                    cell.alignment = Alignment(horizontal="left")

        # Auto-adjust column widths
        for col in ws.columns:
            max_len = max(len(str(cell.value or '')) for cell in col)
            col_letter = get_column_letter(col[0].column)
            ws.column_dimensions[col_letter].width = max(max_len + 4, 12)

        wb.save(str(file_path))

        # Open in default spreadsheet app
        if sys.platform == "win32":
            os.startfile(str(file_path))
        elif sys.platform == "darwin":
            subprocess.run(["open", str(file_path)])
        else:
            subprocess.run(["xdg-open", str(file_path)])

        if player:
            player.write_log(f"[DocGen] 📈 Generated XLSX: {filename}")

        return f"Successfully generated and opened Excel spreadsheet: '{filename}' ({len(rows)} rows).\nSaved to: {file_path}"

    except Exception as e:
        print(f"[DocGen] XLSX error: {e}")
        return f"Failed to generate spreadsheet: {e}"


def document_generator_action(
    parameters: dict = None,
    player=None,
    speak=None,
) -> str:
    """
    parameters:
        type   : 'presentation' | 'spreadsheet' | 'pdf' | 'pptx' | 'xlsx'
        topic  : topic or prompt describing the document
        count  : number of slides / rows (optional)
    """
    params = parameters or {}
    doc_type = params.get("type", "").lower().strip()
    topic = params.get("topic") or params.get("description") or "System Summary"
    count = int(params.get("count", 5))

    if doc_type in ("presentation", "pptx", "powerpoint", "slides"):
        return generate_presentation(topic=topic, num_slides=count, player=player)
    elif doc_type in ("spreadsheet", "xlsx", "excel", "sheet"):
        return generate_spreadsheet(topic=topic, player=player)
    else:
        # Auto-detect from topic
        if any(w in topic.lower() for w in ("presentation", "slide", "ppt")):
            return generate_presentation(topic=topic, num_slides=count, player=player)
        else:
            return generate_spreadsheet(topic=topic, player=player)
